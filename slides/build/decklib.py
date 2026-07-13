# decklib.py — shared design system for BADM 225 lecture decks
# All content original. Fully accessible: title placeholders on every slide,
# real text (no text-in-images), high contrast, alt text, logical add-order.
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette: deep forest green (professional, distinctly not-publisher)
INK      = RGBColor(0x21, 0x2A, 0x26)   # near-black green for body text
GREEN    = RGBColor(0x0E, 0x4D, 0x33)   # deep forest — dominant
MOSS     = RGBColor(0x5C, 0x8A, 0x6E)   # supporting tone
TINT     = RGBColor(0xEA, 0xF1, 0xEB)   # pale green card background
GOLD     = RGBColor(0xC9, 0x93, 0x1A)   # sharp accent — large/bold text on dark bg only
GOLD_TEXT= RGBColor(0x8A, 0x64, 0x10)   # darker gold: passes 4.5:1 on white for small text
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
MUTED    = RGBColor(0x5E, 0x6B, 0x64)

HEAD = "Cambria"   # safe-list serif for headers
BODY = "Calibri"   # safe-list sans for body

SW, SH = Inches(13.333), Inches(7.5)

def new_deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = SW, SH
    return prs

def _alt(shape, desc):
    cnv = shape._element.xpath('.//*[local-name()="cNvPr"]')[0]
    cnv.set('descr', desc)

def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def _footer(slide, deck_label, idx):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(6), Inches(0.35))
    tf = tb.text_frame; tf.word_wrap = False
    r = tf.paragraphs[0].add_run(); r.text = deck_label
    r.font.name = BODY; r.font.size = Pt(10); r.font.color.rgb = MUTED
    _alt(tb, "Footer: " + deck_label)
    tb2 = slide.shapes.add_textbox(Inches(12.3), Inches(7.05), Inches(0.7), Inches(0.35))
    tf2 = tb2.text_frame
    r2 = tf2.paragraphs[0].add_run(); r2.text = str(idx)
    r2.font.name = BODY; r2.font.size = Pt(10); r2.font.color.rgb = MUTED
    tf2.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _alt(tb2, "Slide number " + str(idx))

def _title_ph(prs, slide, text, color, size=34, top=Inches(0.45), left=Inches(0.55), width=None):
    # use the layout's title placeholder when present for accessibility
    ph = None
    for shp in slide.placeholders:
        if shp.placeholder_format.idx == 0:
            ph = shp; break
    if ph is None:
        ph = slide.shapes.add_textbox(left, top, width or Inches(12.2), Inches(1.0))
    else:
        ph.left, ph.top, ph.width, ph.height = left, top, width or Inches(12.2), Inches(1.0)
    tf = ph.text_frame; tf.word_wrap = True
    tf.paragraphs[0].text = ""
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    r = tf.paragraphs[0].add_run(); r.text = text
    r.font.name = HEAD; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = color
    _alt(ph, "Slide title: " + text)
    return ph

def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text

def para(tf, text, size=16, color=INK, bold=False, bullet=False, lead=None,
         space_after=8, first=False, align=PP_ALIGN.LEFT, font=BODY, italic=False):
    p = tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    if bullet:
        rb = p.add_run(); rb.text = "▪  "
        rb.font.name = BODY; rb.font.size = Pt(size); rb.font.color.rgb = GOLD_TEXT
    if lead:
        rl = p.add_run(); rl.text = lead + "  "
        rl.font.name = font; rl.font.size = Pt(size); rl.font.bold = True; rl.font.color.rgb = GREEN
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.italic = italic
    return p

def textbox(slide, x, y, w, h, alt="Text"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tb.text_frame.word_wrap = True
    _alt(tb, alt)
    return tb.text_frame

def card(slide, x, y, w, h, fill=TINT, line=None, alt="Content card"):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    _alt(sh, alt)
    tf = sh.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.14); tf.margin_bottom = Inches(0.12)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    return sh, tf

def badge(slide, x, y, d, symbol, fill=GREEN, fg=WHITE, alt_text="Icon"):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = symbol
    r.font.size = Pt(int(d / Inches(1) * 26)); r.font.bold = True
    r.font.color.rgb = fg; r.font.name = BODY
    _alt(sh, alt_text)
    return sh

# ---------- slide constructors ----------
def title_slide(prs, kicker, title, subtitle, deck_label):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _bg(slide, GREEN)
    tfk = textbox(slide, Inches(0.9), Inches(1.5), Inches(11.5), Inches(0.5), "Course label")
    para(tfk, kicker, size=15, color=RGBColor(0xBF, 0xD8, 0xC6), first=True, font=BODY)
    ph = _title_ph(prs, slide, title, WHITE, size=48, top=Inches(2.1), left=Inches(0.9), width=Inches(11.5))
    ph.height = Inches(1.9)
    tfs = textbox(slide, Inches(0.9), Inches(4.2), Inches(11.0), Inches(1.0), "Subtitle")
    para(tfs, subtitle, size=19, color=TINT, first=True)
    tfd = textbox(slide, Inches(0.9), Inches(6.4), Inches(11.0), Inches(0.6), "Attribution")
    para(tfd, deck_label + "  ·  Original instructional material — no publisher content", size=12,
         color=RGBColor(0xA8, 0xC4, 0xB2), first=True)
    return slide

def section_slide(prs, number, title, note_text, deck_label, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _bg(slide, GREEN)
    tfn = textbox(slide, Inches(0.9), Inches(1.3), Inches(3.0), Inches(1.6), "Section number " + number)
    para(tfn, number, size=80, color=GOLD, bold=True, first=True, font=HEAD)
    _title_ph(prs, slide, title, WHITE, size=40, top=Inches(3.1), left=Inches(0.9), width=Inches(11.5))
    tfs = textbox(slide, Inches(0.9), Inches(4.4), Inches(10.5), Inches(1.2), "Section note")
    para(tfs, note_text, size=17, color=TINT, first=True)
    _footer(slide, deck_label, idx)
    return slide

def content_slide(prs, title, deck_label, idx, title_size=32):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _bg(slide, WHITE)
    _title_ph(prs, slide, title, GREEN, size=title_size)
    _footer(slide, deck_label, idx)
    return slide

def bullets_slide(prs, title, items, deck_label, idx, size=17, kicker=None):
    slide = content_slide(prs, title, deck_label, idx)
    tf = textbox(slide, Inches(0.6), Inches(1.7), Inches(12.1), Inches(5.1), "Key points")
    if kicker:
        para(tf, kicker, size=15, color=MUTED, first=True, italic=True, space_after=14)
    for it in items:
        lead, rest = (it if isinstance(it, tuple) else (None, it))
        para(tf, rest, size=size, bullet=True, lead=lead, space_after=12,
             first=(not kicker and it is items[0]))
    return slide

def two_col_slide(prs, title, left_head, left_items, right_head, right_items,
                  deck_label, idx, size=15, left_fill=TINT, right_fill=TINT,
                  left_head_color=GREEN, right_head_color=GREEN):
    slide = content_slide(prs, title, deck_label, idx)
    for (hx, head, items, fill, hc) in (
        (Inches(0.6), left_head, left_items, left_fill, left_head_color),
        (Inches(6.85), right_head, right_items, right_fill, right_head_color)):
        _, tf = card(slide, hx, Inches(1.7), Inches(5.9), Inches(5.0), fill=fill,
                     alt=head + " panel")
        para(tf, head, size=18, bold=True, color=hc, first=True, space_after=12, font=HEAD)
        for it in items:
            lead, rest = (it if isinstance(it, tuple) else (None, it))
            para(tf, rest, size=size, bullet=True, lead=lead, space_after=9)
    return slide

def icon_rows_slide(prs, title, rows, deck_label, idx, size=15.5, kicker=None):
    # rows: list of (symbol, header, description)
    slide = content_slide(prs, title, deck_label, idx)
    y = Inches(1.75)
    if kicker:
        tf = textbox(slide, Inches(0.6), Inches(1.55), Inches(12.0), Inches(0.45), "Context")
        para(tf, kicker, size=14, color=MUTED, first=True, italic=True)
        y = Inches(2.1)
    n = len(rows)
    row_h = min(Inches(1.18), Inches(5.0) / n)
    for (sym, head, desc) in rows:
        badge(slide, Inches(0.7), y + Inches(0.06), Inches(0.62), sym,
              alt_text=head + " icon")
        tf = textbox(slide, Inches(1.55), y, Inches(11.2), row_h, head)
        para(tf, head, size=size + 1, bold=True, color=GREEN, first=True, space_after=2, font=HEAD)
        para(tf, desc, size=size, space_after=0)
        y += row_h + Inches(0.12)
    return slide

def flow_slide(prs, title, steps, deck_label, idx, note_text=None, size=14):
    # steps: list of (label, sub) — horizontal chevron process
    slide = content_slide(prs, title, deck_label, idx)
    n = len(steps)
    gap = Inches(0.12)
    w = (Inches(12.2) - gap * (n - 1)) / n
    x = Inches(0.6)
    for i, (label, sub) in enumerate(steps):
        sh = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, x, Inches(2.0), w, Inches(1.25))
        sh.fill.solid(); sh.fill.fore_color.rgb = GREEN if i % 2 == 0 else MOSS
        sh.line.fill.background(); sh.shadow.inherit = False
        _alt(sh, "Step " + str(i + 1) + ": " + label)
        tf = sh.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = label
        r.font.name = BODY; r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = WHITE
        tfs = textbox(slide, x, Inches(3.45), w, Inches(1.6), label + " detail")
        para(tfs, sub, size=12.5, color=INK, first=True, align=PP_ALIGN.CENTER)
        x += w + gap
    if note_text:
        _, tf = card(slide, Inches(0.6), Inches(5.35), Inches(12.1), Inches(1.35),
                     alt="Key insight")
        para(tf, note_text, size=15, first=True, lead="Key insight:")
    return slide

def stat_slide(prs, title, big, big_label, points, deck_label, idx):
    slide = content_slide(prs, title, deck_label, idx)
    _, tfL = card(slide, Inches(0.6), Inches(1.8), Inches(4.6), Inches(4.6),
                  fill=GREEN, alt="Highlighted statistic")
    p = tfL.paragraphs[0]; p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(4)
    r = p.add_run(); r.text = big
    r.font.name = HEAD; r.font.size = Pt(64); r.font.bold = True; r.font.color.rgb = WHITE
    para(tfL, big_label, size=15, color=TINT, align=PP_ALIGN.CENTER)
    tfR = textbox(slide, Inches(5.7), Inches(1.9), Inches(7.0), Inches(4.7), "Supporting points")
    for i, it in enumerate(points):
        lead, rest = (it if isinstance(it, tuple) else (None, it))
        para(tfR, rest, size=16, bullet=True, lead=lead, space_after=14, first=(i == 0))
    return slide

def quote_slide(prs, title, quote, attribution, deck_label, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _bg(slide, TINT)
    _title_ph(prs, slide, title, GREEN)
    tf = textbox(slide, Inches(1.4), Inches(2.4), Inches(10.5), Inches(2.6), "Quotation")
    para(tf, "“" + quote + "”", size=27, color=INK, first=True, font=HEAD, italic=True)
    tfa = textbox(slide, Inches(1.4), Inches(5.2), Inches(10.5), Inches(0.6), "Attribution")
    para(tfa, "— " + attribution, size=16, color=MUTED, first=True)
    _footer(slide, deck_label, idx)
    return slide

def discussion_slide(prs, title, questions, deck_label, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _bg(slide, GREEN)
    _title_ph(prs, slide, title, WHITE)
    tf = textbox(slide, Inches(0.8), Inches(1.9), Inches(11.7), Inches(4.9), "Discussion questions")
    for i, q in enumerate(questions):
        para(tf, q, size=18, color=WHITE, space_after=18, first=(i == 0),
             lead=str(i + 1) + ".")
        tf.paragraphs[-1].runs[0].font.color.rgb = GOLD
    _footer_dark(slide, deck_label, idx)
    return slide

def _footer_dark(slide, deck_label, idx):
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(6), Inches(0.35))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = deck_label
    r.font.name = BODY; r.font.size = Pt(10); r.font.color.rgb = RGBColor(0xA8, 0xC4, 0xB2)
    _alt(tb, "Footer")

def takeaways_slide(prs, items, deck_label, idx, site_note=None):
    slide = content_slide(prs, "Key takeaways", deck_label, idx)
    y = Inches(1.7)
    step = Inches(0.7) if site_note else Inches(0.8)
    for i, it in enumerate(items):
        badge(slide, Inches(0.7), y + Inches(0.02), Inches(0.46), "✓",
              alt_text="Checkmark")
        tf = textbox(slide, Inches(1.45), y, Inches(11.3), Inches(0.7), "Takeaway " + str(i + 1))
        para(tf, it, size=15.5, first=True)
        y += step
    if site_note:
        cy = min(y + Inches(0.08), Inches(6.15))
        _, tf = card(slide, Inches(0.6), cy, Inches(11.4), Inches(0.75),
                     fill=GREEN, alt="Practice reminder")
        para(tf, site_note, size=14.5, color=WHITE, first=True, space_after=0)
        tf.paragraphs[0].runs[0].font.bold = True
    return slide
